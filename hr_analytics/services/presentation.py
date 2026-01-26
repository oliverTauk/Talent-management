from __future__ import annotations
from typing import Dict, Any, Optional

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None  # type: ignore


class PresentationService:
    """
    Minimal scaffolding to populate a PPTX from a template.
    Expect a template with text placeholders; replace by key mapping.
    """

    def __init__(self) -> None:
        pass

    def populate_from_template(self, template_path: str, replacements: Dict[str, str], output_path: str) -> Optional[str]:
        if Presentation is None:
            return None
        prs = Presentation(template_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, 'has_text_frame'):
                    continue
                if not shape.has_text_frame:
                    continue
                text = shape.text
                for k, v in replacements.items():
                    if k in text:
                        shape.text = text.replace(k, v)
        prs.save(output_path)
        return output_path
